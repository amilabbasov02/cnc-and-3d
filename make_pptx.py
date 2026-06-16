# -*- coding: utf-8 -*-
"""FEEDRATE — İnvestor pitch deck (.pptx). Premium launch + mövqeləşmə."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import model as M

INK   = RGBColor(0x0D,0x0E,0x10)
SURF  = RGBColor(0x17,0x1A,0x1E)
AMBER = RGBColor(0xFF,0xB3,0x00)
CYAN  = RGBColor(0x46,0xC9,0xC1)
TEXT  = RGBColor(0xEC,0xEA,0xE5)
DIM   = RGBColor(0x9A,0x9E,0xA4)
GREEN = RGBColor(0x57,0xC9,0x8A)
RED   = RGBColor(0xFF,0x6B,0x6B)
WHITE = RGBColor(0xFF,0xFF,0xFF)

rows,_ = M.compute(months=48)
rows0,_ = M.compute({"capital":0}, months=48)
mincum0 = min(r['cum'] for r in rows0)
need = -mincum0
raise_amt = M.A["capital"]
def at(m): return rows[m-1]
def k(x): return f"${x/1000:,.0f}K"
def mm(x): return f"${x/1_000_000:.2f}M"
yr={1:[0,0,0],2:[0,0,0],3:[0,0,0],4:[0,0,0]}
for r in rows:
    g=1 if r['m']<=12 else 2 if r['m']<=24 else 3 if r['m']<=36 else 4
    yr[g][0]+=r['rev']; yr[g][1]+=r['cost']; yr[g][2]+=r['net']
breakeven=next((r['m'] for r in rows if r['net']>0), None)

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def slide(bg=INK):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(1,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=bg
    r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s
def box(s,l,t,w,h):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tb.text_frame.word_wrap=True; return tb.text_frame
def para(tf,text,size,color,bold=False,first=False,align=PP_ALIGN.LEFT,space=6,italic=False):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(space)
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color; f.name="Calibri"; return p
def accent_bar(s):
    b=s.shapes.add_shape(1,0,0,Inches(0.18),Inches(7.5)); b.fill.solid()
    b.fill.fore_color.rgb=AMBER; b.line.fill.background(); b.shadow.inherit=False
def kicker(s,t): tf=box(s,0.7,0.5,12,0.5); para(tf,t,13,AMBER,bold=True,first=True)
def title(s,t,size=34,top=0.95): tf=box(s,0.7,top,12,1.2); para(tf,t,size,WHITE,bold=True,first=True)
def bullets(s,items,top=2.2,size=18,left=0.8,width=11.8):
    tf=box(s,left,top,width,7.5-top-0.4)
    for i,(txt,*opt) in enumerate(items):
        c=opt[0] if opt else TEXT; b=opt[1] if len(opt)>1 else False
        para(tf,("•  " if not b else "")+txt,size,c,bold=b,first=(i==0),space=9)
def stat_card(s,l,t,w,h,big,label,bigcolor=AMBER):
    c=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb=SURF; c.line.color.rgb=AMBER; c.line.width=Pt(1); c.shadow.inherit=False
    tf=c.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    para(tf,big,28,bigcolor,bold=True,first=True,align=PP_ALIGN.CENTER,space=2)
    para(tf,label,12,DIM,align=PP_ALIGN.CENTER,space=0)

# 1 Title
s=slide(); accent_bar(s)
tf=box(s,0.9,2.4,11.5,2.6)
para(tf,"FEEDRATE",60,WHITE,bold=True,first=True,space=4)
para(tf,"CNC frezeleme & 3D-çap üçün anında qiymət → sifariş → istehsal → çatdırılma platforması",22,AMBER,space=18)
para(tf,"Self-serve · brauzer CAD · supplier şəbəkəsi · marketplace · orta bazardan yuxarıya",16,DIM)
tf2=box(s,0.9,6.2,11,0.8); para(tf2,"İnvestor təqdimatı  ·  Kayzen  ·  support@kayzen.az  ·  2026",13,DIM,first=True)

# 2 Problem
s=slide(); accent_bar(s); kicker(s,"PROBLEM"); title(s,"Hissə istehsal etdirmək yavaş, bahalı və qırıqdır")
bullets(s,[
 ("CNC/3D sifarişdə qiymət almaq günlər çəkir — email, gözləmə, geriyə-irəli yazışma",),
 ("Mövcud iri platformalar yalnız böyük sifarişlərə baxır — minimum həcm, sales-led satış",),
 ("Kiçik biznes, mühəndis, startap, maker — sürətli, sərfəli, self-serve həll tapmır",),
 ("Qiymət ayrı, CAD ayrı, sifariş ayrı, emalatxana ayrı — vahid axın yoxdur",),
 ("Nəticə: orta bazar xidmətsiz qalır — bizim fürsətimiz məhz buradadır",AMBER,True),
])

# 3 Həll
s=slide(); accent_bar(s); kicker(s,"HƏLL"); title(s,"Fayldan hazır hissəyə — bir platformada")
bullets(s,[
 ("Anında qiymət: model yüklə / kitabxanadan seç / brauzer CAD-da qur → saniyələrlə qiymət",),
 ("Tək kliklə sifariş + ödəniş — minimum həcm yoxdur, yazışma yoxdur, gözləmə yoxdur",),
 ("Supplier şəbəkəsi: sifariş ən uyğun emalatxanaya yönlənir, hazırlanır, çatdırılır",),
 ("Biz istehsalı etmirik — qiymət, ödəniş, marşrutlaşdırma, keyfiyyəti idarə edirik (yüngül model)",),
 ("Üstəlik: dizayn marketplace — istifadəçilər öz modellərini satır, biz komissiya alırıq",GREEN,True),
])

# 4 Necə işləyir
s=slide(); accent_bar(s); kicker(s,"NECƏ İŞLƏYİR"); title(s,"Qiymətdən çatdırılmaya — 7 addım")
steps=["1 · Model (yüklə / seç / CAD-da qur)","2 · Anında qiymət","3 · Sifariş + ödəniş",
 "4 · Supplier-ə marşrutlaşdırma","5 · İstehsal (CNC/3D-print)","6 · Keyfiyyət + çatdırılma","7 · Supplier payout · biz marja saxlayırıq"]
tf=box(s,0.9,2.2,11.6,4.5)
for i,st in enumerate(steps):
    para(tf,st,19,(AMBER if i in (1,6) else TEXT),bold=(i in(2,6)),first=(i==0),space=11)
tf2=box(s,0.9,6.6,12,0.7); para(tf2,"İstehsalı biz etmirik — qiymət, ödəniş, marşrutlaşdırma və keyfiyyəti biz idarə edirik.",14,CYAN,italic=True,first=True)

# 5 Məhsul
s=slide(); accent_bar(s); kicker(s,"MƏHSUL"); title(s,"Üç modul, bir axın")
bullets(s,[
 ("3D model kitabxanası — hazır parametrik modellər, ölçü qoy → qiymət",),
 ("Brauzer CAD redaktoru (Three.js + CSG) — quraşdırma yox, anlıq qiymət, STL/OBJ/STEP export",),
 ("İstehsal marketplace — sifariş → supplier → keyfiyyət → çatdırılma + payout (Stripe Connect)",),
 ("Dizayn marketplace — listing, alış, komissiya, payout",),
 ("Limit/anti-abuse: anonim 1 qiymət → pulsuz hesab 5/ay → ödənişli paketlər",DIM,False),
])

# 6 Bazar
s=slide(); accent_bar(s); kicker(s,"BAZAR"); title(s,"Böyüyən on-demand istehsal bazarı")
stat_card(s,0.9,2.4,3.7,2.0,"TAM","Qlobal rəqəmsal istehsal —\nonlarla milyard $")
stat_card(s,4.85,2.4,3.7,2.0,"SAM","Self-serve instant-quote +\nonlayn CAD istifadəçiləri")
stat_card(s,8.8,2.4,3.6,2.0,"SOM","4 ildə aylıq 300K+ ziyarətçi,\n5700+ sifariş/ay",GREEN)
tf=box(s,0.9,4.9,11.6,2)
para(tf,"Trendlər: distributed manufacturing, 3D-print bumu, brauzer/no-code alətlər, qısa tədarük zənciri.",16,TEXT,first=True,space=8)
para(tf,"Giriş nöqtəsi: anonim 1 pulsuz qiymət = viral, sürtünməsiz cəlb.",16,AMBER)

# 7 Biznes modeli
s=slide(); accent_bar(s); kicker(s,"BİZNES MODELİ"); title(s,"3 gəlir mənbəyi")
stat_card(s,0.9,2.4,3.7,2.3,"~22%","İSTEHSAL MARJASI (əsas)\nhər sifariş GMV-dən")
stat_card(s,4.85,2.4,3.7,2.3,"$9–$99","ABUNƏ (SaaS)\nCAD + qiymət alətləri / ay")
stat_card(s,8.8,2.4,3.6,2.3,"~8%","DİZAYN KOMİSSİYA\nmarketplace satışından")
tf=box(s,0.9,5.1,11.6,2)
para(tf,"İstehsal: ödənişi biz toplayırıq, supplier-ə payout edirik, marjanı saxlayırıq.",15,TEXT,first=True,space=6)
para(tf,"Üç axın = həcmlə böyüyən marja + təkrar MRR + iki tərəfli marketplace effekti.",15,GREEN)

# 8 Supplier
s=slide(); accent_bar(s); kicker(s,"SUPPLIER ŞƏBƏKƏSİ"); title(s,"İstehsalsız, lakin tam nəzarətli")
bullets(s,[
 ("Yoxlanmış CNC/3D-print emalatxanaları şəbəkəsi — biz onlara sifariş gətiririk",),
 ("Marşrutlaşdırma: hər sifariş qiymət/məkan/imkana görə ən uyğun supplier-ə",),
 ("Keyfiyyət standartı, izləmə və reytinq — zəif supplier sistemdən çıxır",),
 ("Bizdə əsas vəsait (avadanlıq) yoxdur → yüngül, sürətli miqyaslana bilən model",),
 ("Çox supplier = aşağı qiymət, sürətli çatdırılma, regional əhatə",GREEN,True),
])

# 9 Qiymət
s=slide(); accent_bar(s); kicker(s,"QİYMƏT"); title(s,"Freemium → ödənişli paketlər")
bullets(s,[
 ("Free $0 — 5 hesablama / 5 yükləmə, məhdud CAD, sifariş icazəli",),
 ("Starter $9 — 50 hesablama / 25 yükləmə, tam CAD",),
 ("Pro $29 — limitsiz hesablama / 100 yükləmə, brendli PDF",),
 ("Business $99 — limitsiz, API/komanda, prioritet supplier & dəstək (yuxarı bazar)",),
 ("Blended ARPU ~$26/ay · istehsal sifarişləri paketdən asılı deyil — hamı sifariş verə bilər",AMBER,True),
])

# 10 Rəqabət
s=slide(); accent_bar(s); kicker(s,"RƏQABƏT"); title(s,"Rəqiblər kimdir")
bullets(s,[
 ("Xometry / Protolabs — instant-quote var, AMMA enterprise/yüksək bazar, sales-led, minimum həcm",),
 ("Protolabs öz zavodları ilə premium qiymət — kiçik/orta sifarişçi üçün baha və əlçatmaz",),
 ("Thangs / Thingiverse — model var, qiymət/sifariş yox",),
 ("Onshape / Tinkercad — CAD var, qiymət + istehsal + marketplace yox",),
 ("Heç biri orta bazara self-serve, hamısını-bir-yerdə həll vermir — boşluq buradadır",AMBER,True),
])

# 11 Bazar mövqeyi — NİYƏ ORTA BAZAR (disruption from below)
s=slide(); accent_bar(s); kicker(s,"MÖVQELƏŞMƏ"); title(s,"Niyə orta bazar? — aşağıdan yuxarıya disrupsiya")
bullets(s,[
 ("Rəqiblər yüksək/enterprise bazara fokuslanır — iri müştəri, böyük marja, sales komandası",),
 ("Orta bazar (kiçik biznes, mühəndis, startap, hardware maker, tələbə) onlar üçün «kiçik»dir — xidmətsiz qalır",),
 ("Bu seqment böyükdür, sürətli böyüyür və self-serve istəyir — bizim giriş nöqtəmiz",AMBER,True),
 ("Strategiya: orta bazarı self-serve ilə tut → həcm və data qazan → Business/API ilə yuxarı bazara çıx",),
 ("Klassik «disruption from below»: əvvəl çatışmayan seqment, sonra premium müştəriyə doğru hərəkət",GREEN,True),
 ("Aşağı sürtünmə + minimum həcm yoxdur = rəqiblərin əlçatmaz olduğu kütləvi tələbi açırıq",),
])

# 12 Müqayisə cədvəli — bizdə var / rəqiblərdə yox
s=slide(); accent_bar(s); kicker(s,"FƏRQLƏNMƏ"); title(s,"Bizdə var — rəqiblərdə yox")
data=[
 ["Funksiya","FEEDRATE","Xometry / Protolabs","Onshape / Thingiverse"],
 ["Self-serve anında qiymət","✓","Qismən (iri/B2B)","✗"],
 ["Brauzer CAD (quraşdırmasız)","✓","✗","Qismən (yalnız CAD)"],
 ["Sifariş → supplier istehsal","✓","✓","✗"],
 ["Minimum həcm YOXDUR","✓","✗","—"],
 ["Dizayn marketplace","✓","✗","Qismən"],
 ["Hamısı bir platformada","✓","✗","✗"],
 ["Orta bazara fokus","✓","✗ (enterprise)","—"],
]
rowsN=len(data); colsN=4
tbl=s.shapes.add_table(rowsN,colsN,Inches(0.7),Inches(2.05),Inches(11.9),Inches(4.9)).table
tbl.columns[0].width=Inches(4.0); tbl.columns[1].width=Inches(2.5)
tbl.columns[2].width=Inches(2.9); tbl.columns[3].width=Inches(2.5)
for ci in range(colsN):
    for ri in range(rowsN):
        cell=tbl.cell(ri,ci); cell.margin_top=Pt(2); cell.margin_bottom=Pt(2)
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        val=data[ri][ci]
        if ri==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=AMBER; col=INK; b=True; sz=13
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb=(SURF if ri%2 else RGBColor(0x20,0x24,0x28))
            b=(ci==1); sz=12.5
            if ci==1 and val=="✓": col=GREEN
            elif val=="✗": col=RED
            else: col=TEXT
        cell.text=val
        p=cell.text_frame.paragraphs[0]; p.alignment=(PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
        r=p.runs[0]; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=col; r.font.name="Calibri"

# 13 Get-to-market
s=slide(); accent_bar(s); kicker(s,"GET-TO-MARKET"); title(s,"Necə böyüyəcəyik")
bullets(s,[
 ("SEO + kontent (CNC/3D-print qiymət açar sözləri) — anonim 1 pulsuz qiymət viral giriş",),
 ("Performans reklam + maker icmaları: Reddit, Discord, YouTube, mühəndis forumları",),
 ("Marketplace satıcıları öz auditoriyasını gətirir — iki tərəfli böyümə",),
 ("Supplier tərəfdaşlıqları → coğrafi genişlənmə",),
 ("Güclü launch kampaniyası (premium büdcə) + daimi CAC/LTV optimallaşdırması",AMBER,True),
])

# 14 Yol xəritəsi
s=slide(); accent_bar(s); kicker(s,"YOL XƏRİTƏSİ"); title(s,"İlk 2 ay intensiv build — sonra güclü launch")
bullets(s,[
 ("Ay 1–2: İNTENSİV BUILD (gəlir yoxdur) — 18 nəfərlik komanda ilə Auth, DB, qiymət mühərriki, CAD, ödəniş, supplier-lər",AMBER,True),
 ("Ay 3: LAUNCH — güclü marketinq, ilk sifarişlər, supplier istehsalı, gəlir başlayır",GREEN,True),
 ("Ay 3–10: Brauzer CAD v1, marketplace, tam sifariş→çatdırılma dövrü",),
 ("Ay 14–24: Orta bazarda böyümə, supplier şəbəkəsinin genişlənməsi",),
 ("Ay 24–36: Business/API ilə yuxarı bazara çıxış, STEP export, beynəlxalq genişlənmə, mənfəət",),
])

# 15 Komanda
s=slide(); accent_bar(s); kicker(s,"KOMANDA"); title(s,"18 nəfərlik intensiv build komandası")
bullets(s,[
 ("Ay 1 nüvə: CEO + CTO + HR + 2 Senior CAD + Frontend + Backend + UI/UX + DevOps",),
 ("Ay 2 tamamlanma: +Frontend, +Backend, +UI/UX, 2 QA, CAD QA, Dizayn QA, Ops, Marketinq (18 nəfər)",),
 ("Ay 6–12: 3D kontent, CFO, Sales/BizDev, Müştəri dəstəyi",),
 ("Ay 22–30: əlavə dəstək və developer (miqyas)",),
 (f"36 ayın sonunda {M.headcount(36)} nəfərlik güclü komanda · remote-first · Azərbaycan bazlı",DIM,False),
])

# 16 Maliyyə
s=slide(); accent_bar(s); kicker(s,"MALİYYƏ"); title(s,"Harda olacağıq — 18, 36, 48 ay")
m18,m36,m48=at(18),at(36),at(48)
stat_card(s,0.9,2.3,3.7,1.55,k(m18['rev']),"18 AY — aylıq gəlir")
stat_card(s,4.85,2.3,3.7,1.55,k(m36['rev']),"36 AY — aylıq gəlir")
stat_card(s,8.8,2.3,3.6,1.55,k(m48['rev']),"48 AY — aylıq gəlir",GREEN)
stat_card(s,0.9,4.05,3.7,1.55,k(m18['rev']*12),"18 AY — ARR")
stat_card(s,4.85,4.05,3.7,1.55,mm(m36['rev']*12),"36 AY — ARR")
stat_card(s,8.8,4.05,3.6,1.55,mm(m48['rev']*12),"48 AY — ARR",GREEN)
tf=box(s,0.9,5.85,11.6,1.4)
para(tf,f"İl 1 {k(yr[1][0])} · İl 2 {k(yr[2][0])} · İl 3 {mm(yr[3][0])} · İl 4 {mm(yr[4][0])} gəlir   |   Breakeven ≈ {breakeven}-ci ay",15,TEXT,first=True,space=5)
para(tf,f"48 ayda sifariş {m48['orders']:,.0f}/ay · ödənişli abunəçi {m48['paid']:,.0f} · komanda {m48['head']} nəfər · İl 4 NET {mm(yr[4][2])}",14,DIM)

# 17 İnvestisiya
s=slide(); accent_bar(s); kicker(s,"İNVESTİSİYA TƏLƏBİ"); title(s,f"Tələb: {mm(raise_amt)} seed")
stat_card(s,0.9,2.3,3.7,1.7,mm(raise_amt),"Seed investisiya (köklü start)")
stat_card(s,4.85,2.3,3.7,1.7,"~30 ay","Runway (breakeven-ə qədər)")
stat_card(s,8.8,2.3,3.6,1.7,mm(m48['rev']*12),"48 ay ARR hədəfi",GREEN)
bullets(s,[
 ("İstifadə: komanda ~60% · məhsul/infra ~12% · marketinq ~20% · əməliyyat/inzibati ~8%",),
 ("İlk 2 ay intensiv build (gəlirsiz) tam maliyyələşdirilir → ay 3-də güclü launch",),
 (f"Plan: orta bazarı tut → {mm(m48['rev']*12)} ARR-ə çat → yuxarı bazara genişlən, İl 4 NET {mm(yr[4][2])}",GREEN,True),
],top=4.3,size=16)

# 18 Qapanış
s=slide(); accent_bar(s)
tf=box(s,0.9,2.7,11.5,2.5)
para(tf,"Gəlin birlikdə quraq.",44,WHITE,bold=True,first=True,space=14)
para(tf,"FEEDRATE — fayldan hazır hissəyə, anında. Orta bazardan başlayıb yuxarıya.",22,AMBER,space=24)
para(tf,"Kayzen  ·  support@kayzen.az",16,DIM)

prs.save("FEEDRATE-Pitch-Deck-48ay.pptx")
print("OK PPTX · slaydlar:",len(prs.slides._sldIdLst))
print(f"raise={raise_amt:,.0f} need={need:,.0f} 18MRR={m18['rev']:,.0f} 36MRR={m36['rev']:,.0f} 36ARR={m36['rev']*12:,.0f} breakeven={breakeven}")
