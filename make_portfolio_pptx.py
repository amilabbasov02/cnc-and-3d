# -*- coding: utf-8 -*-
"""KAYZEN portfel — investor pitch deck (.pptx). Tək şirkət, axan komanda, 4 məhsul."""
import io, sys
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import portfolio as P, model as FR

INK=RGBColor(0x0D,0x0E,0x10); SURF=RGBColor(0x17,0x1A,0x1E); AMBER=RGBColor(0xFF,0xB3,0x00)
CYAN=RGBColor(0x46,0xC9,0xC1); TEXT=RGBColor(0xEC,0xEA,0xE5); DIM=RGBColor(0x9A,0x9E,0xA4)
GREEN=RGBColor(0x57,0xC9,0x8A); RED=RGBColor(0xFF,0x6B,0x6B); WHITE=RGBColor(0xFF,0xFF,0xFF)
PCLR={"FEEDRATE":RGBColor(0xFF,0xB3,0x00),"QUOTEFLOW":RGBColor(0x5C,0x9C,0xFF),
      "FORMCHECK":RGBColor(0x57,0xC9,0x8A),"CONFIGFLOW":RGBColor(0xB9,0x7D,0xF0)}

rows,meta=P.compute(); rows0,_=P.compute(capital=0); trough0=min(r['cum'] for r in rows0)
need=-trough0; raise_amt=meta['capital']
def at(m): return rows[m-1]
def k(x): return f"${x/1000:,.0f}K"
def mm(x): return f"${x/1_000_000:.2f}M"
yr={1:[0,0,0],2:[0,0,0],3:[0,0,0],4:[0,0,0]}
for r in rows:
    g=1 if r['m']<=12 else 2 if r['m']<=24 else 3 if r['m']<=36 else 4
    yr[g][0]+=r['rev']; yr[g][1]+=r['cost']; yr[g][2]+=r['net']
breakeven=next((r['m'] for r in rows if r['net']>0),None)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height
def slide(bg=INK):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(1,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=bg
    r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return s
def box(s,l,t,w,h):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tb.text_frame.word_wrap=True; return tb.text_frame
def para(tf,t,size,color,bold=False,first=False,align=PP_ALIGN.LEFT,space=6,italic=False):
    p=tf.paragraphs[0] if first else tf.add_paragraph(); p.alignment=align; p.space_after=Pt(space)
    r=p.add_run(); r.text=t; f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color; f.name="Calibri"; return p
def bar(s): b=s.shapes.add_shape(1,0,0,Inches(0.18),Inches(7.5)); b.fill.solid(); b.fill.fore_color.rgb=AMBER; b.line.fill.background(); b.shadow.inherit=False
def kicker(s,t): tf=box(s,0.7,0.5,12,0.5); para(tf,t,13,AMBER,bold=True,first=True)
def title(s,t,size=32,top=0.95): tf=box(s,0.7,top,12,1.2); para(tf,t,size,WHITE,bold=True,first=True)
def bullets(s,items,top=2.2,size=18,left=0.8,width=11.9):
    tf=box(s,left,top,width,7.5-top-0.4)
    for i,(txt,*o) in enumerate(items):
        c=o[0] if o else TEXT; b=o[1] if len(o)>1 else False
        para(tf,("•  " if not b else "")+txt,size,c,bold=b,first=(i==0),space=9)
def card(s,l,t,w,h,big,label,bigc=AMBER,lc=DIM):
    c=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h)); c.fill.solid(); c.fill.fore_color.rgb=SURF
    c.line.color.rgb=bigc; c.line.width=Pt(1); c.shadow.inherit=False
    tf=c.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    para(tf,big,26,bigc,bold=True,first=True,align=PP_ALIGN.CENTER,space=2)
    para(tf,label,11.5,lc,align=PP_ALIGN.CENTER,space=0)
def rect(s,l,t,w,h,color,line=None):
    r=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h)); r.fill.solid(); r.fill.fore_color.rgb=color
    if line: r.line.color.rgb=line; r.line.width=Pt(0.75)
    else: r.line.fill.background()
    r.shadow.inherit=False; return r

# 1 Title
s=slide(); bar(s)
tf=box(s,0.9,2.2,11.6,3)
para(tf,"KAYZEN",58,WHITE,bold=True,first=True,space=4)
para(tf,"İstehsal-texnologiya şirkəti — bir nüvə, ardıcıl çıxan məhsul portfeli",22,AMBER,space=16)
para(tf,"Brauzer CAD + qiymət mühərriki + ödəniş nüvəsi üzərində 4 məhsul · tək komanda · tək şirkət",15,DIM,space=14)
para(tf,f"İnvestor təqdimatı · Kayzen · support@kayzen.az · 2026",13,DIM)

# 2 Vizyon
s=slide(); bar(s); kicker(s,"VİZYON"); title(s,"Tək məhsul deyil — istehsal-texnologiya mühərriki")
bullets(s,[
 ("İstehsal/CAD sahəsi parçalanıb: qiymət ayrı, CAD ayrı, sifariş ayrı, DFM ayrı alətlərdədir",),
 ("Biz bir texnoloji nüvə qururuq (Three.js CAD + qiymət mühərriki + Stripe + analitika)",),
 ("Bu nüvənin üstündə ardıcıl olaraq 4 məhsul çıxarırıq — hər biri ayrı sayt/brend, bir az fərqli bazar",),
 ("Hamısı bir şirkətdə cəmlənir, eyni komanda ilə — komanda heç vaxt boş qalmır",AMBER,True),
 ("İnvestor bir məhsula yox, çoxlu atışlı istehsal-texnologiya platformasına investisiya edir",GREEN,True),
])

# 3 Problem
s=slide(); bar(s); kicker(s,"PROBLEM"); title(s,"Hər istehsal addımı üçün ayrıca, baha, kobud alət")
bullets(s,[
 ("Sifarişçi: hissə qiyməti almaq günlər çəkir, kiçik sifariş üçün self-serve həll yoxdur",),
 ("Emalatxana: qiymət/sifarişi hələ də email + Excel ilə idarə edir — proqram yoxdur",),
 ("Mühəndis: sürətli istehsal-uyğunluq (DFM) yoxlaması üçün ucuz alət tapmır",),
 ("E-commerce satıcı: custom məhsul üçün 3D konfiqurator + anında qiymət qoya bilmir",),
 ("Hər biri ayrıca böyük bazardır — və hamısı eyni texnoloji nüvəni tələb edir",AMBER,True),
])

# 4 Strategiya
s=slide(); bar(s); kicker(s,"STRATEGİYA"); title(s,"Bir nüvə → ardıcıl GTM → bir komanda")
bullets(s,[
 ("1) Paylaşılan nüvəni bir dəfə qur (CAD, qiymət, ödəniş, fayl, analitika) — ~70% hər məhsulda təkrar işlənir",),
 ("2) Hər məhsul bu nüvənin üstündə fərqli bazara «paketlənmiş» go-to-market-dir",),
 ("3) Ardıcıl çıxar: bir məhsul launch olub stabilləşəndə komanda növbətiyə axır",),
 ("4) Hər məhsul növbətindən əvvəl validasiya olunur — risk pillə-pillə azalır",),
 ("Nəticə: 4 məhsul, ~28 nəfər, tək seed — 4 ayrı şirkətin kəsri qədər kapitalla",GREEN,True),
])

# 5 Paylaşılan nüvə
s=slide(); bar(s); kicker(s,"PAYLAŞILAN NÜVƏ"); title(s,"Bir dəfə qurulan, dəfələrlə işlədilən texnologiya")
card(s,0.9,2.4,3.7,1.9,"Brauzer CAD","Three.js + CSG\ngeometriya/mesh emalı")
card(s,4.85,2.4,3.7,1.9,"Qiymət mühərriki","material+vaxt+marja\nno-code qayda builder")
card(s,8.8,2.4,3.6,1.9,"Ödəniş + Connect","Stripe abunə + payout")
card(s,0.9,4.5,3.7,1.9,"Fayl pipeline","STL/STEP/OBJ\nyükləmə, render, saxlama")
card(s,4.85,4.5,3.7,1.9,"Auth + Analitika","istifadə, konversiya,\ndashboard")
card(s,8.8,4.5,3.6,1.9,"~70%","hər yeni məhsulda\ntəkrar işlənən kod",GREEN)

# 6-9 məhsullar
def product_slide(name, sub, what, role, launch):
    s=slide(); bar(s); kicker(s,f"MƏHSUL · {name}"); title(s,sub)
    bullets(s,what,top=2.2,size=18)
    card(s,9.6,0.7,3.0,1.1,f"Launch: Ay {launch}",name,bigc=PCLR[name])
    return s
product_slide("FEEDRATE","İstehsal marketplace — qiymətdən çatdırılmaya",[
 ("Anında qiymət → sifariş → supplier istehsalı → keyfiyyət → çatdırılma",),
 ("Gəlir: istehsal marjası ~22% + abunə + dizayn komissiya",),
 ("Orta bazar fokusu (Xometry/Protolabs yüksək bazara baxır) — disruption from below",AMBER,True),
 ("Rolu: flaqman + nağd axın + supplier şəbəkəsi (digər məhsulları qidalandırır)",GREEN,True),
],None,FR.A["launch_month"])
product_slide("QUOTEFLOW","Embed quote-widget SaaS — emalatxanalar üçün",[
 ("Emalatxana abunə olur, dizaynı öz saytına uyğunlaşdırır, öz qiymət məntiqini yazır, embed götürür",),
 ("Analitika: kim istifadə etdi, neçə quote, neçəsi sifarişə çevrildi",),
 ("Rəqib: DigiFabster/3YOURMIND — amma köhnə UX, CAD redaktoru yoxdur",),
 ("Gəlir: $140/ay abunə · saf B2B SaaS, fiziki icra yox → yüksək marja",GREEN,True),
],None,9)
product_slide("FORMCHECK","DFM analiz SaaS — istehsal-uyğunluq yoxlaması",[
 ("CAD yüklə → divar qalınlığı, tolerans, undercut, qiyməti artıran amillər + AI tövsiyə",),
 ("Boşluq: SMB üçün yüngül, self-serve DFM aləti demək olar yoxdur (aPriori enterprise/bahalı)",),
 ("Gəlir: freemium → $39/ay · brauzer mesh analizi (nüvə təkrar)",GREEN,True),
],None,16)
product_slide("CONFIGFLOW","E-commerce 3D konfiqurator / CPQ",[
 ("Custom məhsul satıcıları (mebel, lövhə, bracket, korpus) öz mağazasına 3D konfiqurator + anında qiymət qoyur",),
 ("Shopify/Woo app store = hazır paylama kanalı",),
 ("Gəlir: $89/ay abunə + tranzaksiya · Three.js + qiymət nüvəsi təkrar",GREEN,True),
],None,24)

# 10 TIMELINE (Gantt hero)
s=slide(); bar(s); kicker(s,"AXAN KOMANDA"); title(s,"Ardıcıl launch — komanda boş qalmır, işdən çıxarma yox")
chart_l, chart_r, top0 = 2.6, 12.9, 2.3
mw=(chart_r-chart_l)/48.0; rh=0.72
# month gridlines/labels (every 6)
for mm_ in range(0,49,6):
    x=chart_l+mm_*mw
    para(box(s,x-0.25,top0-0.42,0.6,0.3),str(mm_ if mm_>0 else 1),11,DIM,first=True,align=PP_ALIGN.CENTER)
prods=[("FEEDRATE",1,FR.A["launch_month"]),("QUOTEFLOW",6,9),("FORMCHECK",13,16),("CONFIGFLOW",21,24)]
for i,(name,bs,lc) in enumerate(prods):
    y=top0+i*rh
    para(box(s,0.55,y+0.02,2.0,0.5),name,13,PCLR[name],bold=True,first=True)
    # build bar (amber outline)
    bx=chart_l+(bs-1)*mw; bw=(lc-bs)*mw
    if bw>0:
        rb=rect(s,bx,y,bw,0.42,SURF,line=AMBER)
        para(rb.text_frame,"build",10,AMBER,first=True,align=PP_ALIGN.CENTER)
    # live bar
    lx=chart_l+(lc-1)*mw; lw=(49-lc)*mw
    rl=rect(s,lx,y,lw,0.42,PCLR[name]); rl.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    para(rl.text_frame,f"canlı · launch Ay{lc}",10,INK,bold=True,first=True,align=PP_ALIGN.CENTER)
tf=box(s,0.55,top0+4*rh+0.15,12.2,1.4)
para(tf,"Ay 3 FEEDRATE launch → komanda QUOTEFLOW-a axır (Ay 9) · Ay 16 → FORMCHECK · Ay 24 → CONFIGFLOW",15,TEXT,first=True,space=5)
para(tf,"Hər keçiddə nüvə komanda saxlanılır, yalnız GTM/məhsul qatı dəyişir — boş qalma və ya işdən çıxarma yoxdur.",14,GREEN,italic=True)

# 11 Kapital səmərəliliyi
s=slide(); bar(s); kicker(s,"KAPİTAL SƏMƏRƏLİLİYİ"); title(s,"4 məhsul — bir komandanın qiymətinə")
card(s,0.9,2.4,3.7,2.0,f"{at(48)['head']} nəfər","48 ayda 4 məhsul üçün")
card(s,4.85,2.4,3.7,2.0,"~72 nəfər","4 ayrı şirkət olsaydı",bigc=RED,lc=DIM)
card(s,8.8,2.4,3.6,2.0,"~70%","təkrar işlənən nüvə kod",GREEN)
bullets(s,[
 ("Paylaşılan back-office: bir CFO, bir marketinq mühərriki, bir infra, bir analitika",),
 ("Hər yeni məhsulun marjinal xərci aşağı — nüvə artıq var",),
 ("FEEDRATE nağd axını digər məhsulların build-ini maliyyələşdirir",GREEN,True),
],top=4.7,size=17)

# 12 Fokus etirazı
s=slide(); bar(s); kicker(s,"«BU FOKUS İTKİSİ DEYİL?»"); title(s,"Xeyr — bir nüvə, ardıcıl icra")
bullets(s,[
 ("Bu, 4 fərqli biznes deyil — bir texnoloji platformanın 4 go-to-market təbəqəsidir",),
 ("Eyni anda 4 şey qurmuruq: ardıcıl çıxarırıq, biri stabilləşəndə digərinə keçirik",),
 ("Hər məhsul növbətiyə START vermədən əvvəl bazarda validasiya olunur",),
 ("Eyni müştəri seqmenti, eyni satış kanalları, eyni brend ekosistemi",),
 ("Risk azalır (çoxlu atış), fokus qalır (bir nüvə, bir komanda, bir vizyon)",AMBER,True),
])

# 13 Flywheel
s=slide(); bar(s); kicker(s,"SİNERGİYA"); title(s,"Məhsullar bir-birini qidalandırır")
bullets(s,[
 ("QUOTEFLOW emalatxanaları → FEEDRATE marketplace üçün hazır supplier şəbəkəsi",),
 ("FORMCHECK istifadəçiləri → FEEDRATE/QUOTEFLOW-da qiymət alıb sifarişə çevrilir",),
 ("CONFIGFLOW e-commerce satıcıları → FEEDRATE-də istehsal sifariş edir",),
 ("Tək hesab, tək ödəniş, tək data — cross-sell və yüksək LTV",),
 ("Hər yeni məhsul bütün portfelin dəyərini artırır (şəbəkə effekti)",GREEN,True),
])

# 14 Bazar
s=slide(); bar(s); kicker(s,"BAZAR"); title(s,"4 bitişik bazar — onlarla milyard $ TAM")
card(s,0.9,2.5,2.9,2.0,"İstehsal","on-demand manufacturing\n(FEEDRATE)")
card(s,4.0,2.5,2.9,2.0,"Quoting SaaS","emalatxana proqramı\n(QUOTEFLOW)",bigc=PCLR["QUOTEFLOW"])
card(s,7.1,2.5,2.9,2.0,"DFM / CAD","mühəndis alətləri\n(FORMCHECK)",bigc=PCLR["FORMCHECK"])
card(s,10.2,2.5,2.5,2.0,"E-com CPQ","konfiqurator\n(CONFIGFLOW)",bigc=PCLR["CONFIGFLOW"])
tf=box(s,0.9,4.9,11.8,1.5)
para(tf,"Hamısı böyüyən rəqəmsal istehsal trendinin içində · bitişik = eyni alıcı, eyni kanal, çarpaz satış.",15,TEXT,first=True,space=6)

# 15 Maliyyə
s=slide(); bar(s); kicker(s,"MALİYYƏ"); title(s,"Birləşmiş portfel — 18 / 36 / 48 ay")
m18,m36,m48=at(18),at(36),at(48)
card(s,0.9,2.2,3.7,1.5,k(m18['rev']),"18 ay — aylıq gəlir")
card(s,4.85,2.2,3.7,1.5,k(m36['rev']),"36 ay — aylıq gəlir")
card(s,8.8,2.2,3.6,1.5,k(m48['rev']),"48 ay — aylıq gəlir",GREEN)
card(s,0.9,3.85,3.7,1.5,k(m18['rev']*12),"18 ay — ARR")
card(s,4.85,3.85,3.7,1.5,mm(m36['rev']*12),"36 ay — ARR")
card(s,8.8,3.85,3.6,1.5,mm(m48['rev']*12),"48 ay — ARR",GREEN)
tf=box(s,0.9,5.6,11.8,1.6)
para(tf,f"İl gəlir: İl1 {k(yr[1][0])} · İl2 {k(yr[2][0])} · İl3 {mm(yr[3][0])} · İl4 {mm(yr[4][0])}   |   Breakeven ≈ Ay {breakeven}",15,TEXT,first=True,space=5)
para(tf,f"İl 4 NET {mm(yr[4][2])} · 4 məhsulun hamısı canlı · komanda {m48['head']} nəfər",14,DIM)

# 16 İnvestisiya
s=slide(); bar(s); kicker(s,"İNVESTİSİYA TƏLƏBİ"); title(s,f"Tələb: {mm(raise_amt)} seed")
card(s,0.9,2.3,3.7,1.7,mm(raise_amt),"Bütün portfel üçün seed")
card(s,4.85,2.3,3.7,1.7,f"~{abs(int(need/1000))}K","Real funding need (dib)")
card(s,8.8,2.3,3.6,1.7,mm(m48['rev']*12),"48 ay ARR hədəfi",GREEN)
bullets(s,[
 ("İstifadə: komanda ~60% · məhsul/infra ~12% · marketinq (4 GTM) ~20% · inzibati ~8%",),
 ("FEEDRATE ilk 2 ay intensiv build, Ay 3 launch — sonra hər məhsul ardıcıl maliyyələşir",),
 (f"Plan: 4 məhsul, {mm(m48['rev']*12)} birləşmiş ARR, İl 4 mənfəətli, ~28 nəfərlə",GREEN,True),
],top=4.3,size=16)

# 17 Komanda
s=slide(); bar(s); kicker(s,"KOMANDA"); title(s,"Senior nüvə, məhsuldan-məhsula axır")
bullets(s,[
 ("Nüvə (Ay 1): CEO + CTO/Baş arxitekt + Senior CAD + Senior Frontend — paylaşılan platforma",),
 ("Hər məhsul üçün 1 product/FS lead + 1 dev + 1 growth əlavə olunur (nüvəyə söykənir)",),
 ("Paylaşılan: CFO, marketinq lead, DevOps, QA, dəstək, sales — bir dəfə, hamı üçün",),
 (f"48 ayın sonunda {at(48)['head']} nəfər · remote-first · Azərbaycan bazlı",DIM,False),
])

# 18 Qapanış
s=slide(); bar(s)
tf=box(s,0.9,2.7,11.6,2.6)
para(tf,"Bir komanda. Bir nüvə. Dörd bazar.",40,WHITE,bold=True,first=True,space=14)
para(tf,"KAYZEN — istehsal-texnologiyasını bir-bir fəth edən şirkət.",20,AMBER,space=22)
para(tf,"Kayzen · support@kayzen.az",15,DIM)

prs.save("KAYZEN-Portfel-Pitch-v2.pptx")
print("OK PPTX slaydlar:",len(prs.slides._sldIdLst),"| raise",raise_amt,"need",round(need),"48ARR",round(m48['rev']*12))
